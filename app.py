# -*- coding: utf-8 -*-
"""
Universal File-to-Text Converter for Google Colab.

This script provides a set of functions to upload a file to a Colab
environment, convert its content to plain text or Markdown, display a preview,
and offer the full text as a downloadable file.

Supported formats:
- Word (.docx)
- Excel (.xlsx)
- PowerPoint (.pptx)
- HTML (.html)
- ZIP archives (.zip) containing supported files
- Plain Text (.txt)
"""

# ==============================================================================
# 1. SETUP: Install required packages quietly
# ==============================================================================
# The -q flag ensures a quiet installation without verbose output.
!pip install -q python-docx openpyxl python-pptx beautifulsoup4 markdownify

# ==============================================================================
# 2. IMPORTS: Import necessary libraries
# ==============================================================================
import os
import zipfile
import docx
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup
from markdownify import markdownify as md
from google.colab import files
import io

# ==============================================================================
# 3. CORE FUNCTIONS
# ==============================================================================

def upload_file_to_colab():
    """
    Handles the file upload process in Google Colab.

    This function opens a file upload dialog in the browser and waits for the
    user to select a file. It returns the name and content of the uploaded file.

    Returns:
        tuple: A tuple containing (filename, file_content) or (None, None) if no file is uploaded.
    """
    print("Please upload a file to convert...")
    try:
        # `files.upload()` returns a dictionary of {filename: content}
        uploaded = files.upload()

        if not uploaded:
            print("No file selected. Aborting.")
            return None, None

        # Get the first uploaded file's name and content
        file_name = next(iter(uploaded))
        file_content = uploaded[file_name]
        print(f"\nSuccessfully uploaded: '{file_name}'")
        return file_name, io.BytesIO(file_content)
    except Exception as e:
        print(f"An error occurred during file upload: {e}")
        return None, None


def convert_file_to_text(file_name, file_stream):
    """
    A universal converter that processes a file stream and extracts text.

    It identifies the file type based on its extension and uses the appropriate
    library to extract text content. For ZIP archives, it recursively processes
    the files within.

    Args:
        file_name (str): The name of the file (e.g., 'document.docx').
        file_stream (io.BytesIO): A byte stream of the file's content.

    Returns:
        str: The extracted text content.
    """
    # Get the file extension to determine the processing method
    _, file_extension = os.path.splitext(file_name)
    file_extension = file_extension.lower()
    
    print(f"Processing file with extension: {file_extension}")
    extracted_text = ""

    try:
        # --- Handle Word documents ---
        if file_extension == '.docx':
            doc = docx.Document(file_stream)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            extracted_text = '\n'.join(full_text)

        # --- Handle Excel spreadsheets ---
        elif file_extension == '.xlsx':
            workbook = openpyxl.load_workbook(file_stream)
            full_text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    # Join cell values with a tab, filter out None values
                    row_text = '\t'.join([str(cell.value) for cell in row if cell.value is not None])
                    full_text.append(row_text)
            extracted_text = '\n'.join(full_text)

        # --- Handle PowerPoint presentations ---
        elif file_extension == '.pptx':
            presentation = Presentation(file_stream)
            full_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            extracted_text = '\n---\n'.join(full_text) # Separate slides

        # --- Handle HTML files (convert to Markdown) ---
        elif file_extension == '.html':
            soup = BeautifulSoup(file_stream.read().decode('utf-8', 'ignore'), "html.parser")
            # Use markdownify to convert HTML body to Markdown
            extracted_text = md(str(soup.body), heading_style="ATX")

        # --- Handle ZIP archives ---
        elif file_extension == '.zip':
            full_text = []
            with zipfile.ZipFile(file_stream, 'r') as zip_ref:
                for name in zip_ref.namelist():
                    # Avoid processing macOS resource fork files
                    if not name.startswith('__MACOSX/'):
                        with zip_ref.open(name) as member_file:
                            # Create a BytesIO stream for the member file
                            member_stream = io.BytesIO(member_file.read())
                            # Recursively call this function for each file in the zip
                            print(f"  -> Processing '{name}' from ZIP archive...")
                            member_text = convert_file_to_text(name, member_stream)
                            full_text.append(f"--- Content from: {name} ---\n{member_text}")
            extracted_text = "\n\n".join(full_text)
            
        # --- Handle plain text files ---
        elif file_extension == '.txt':
            extracted_text = file_stream.read().decode('utf-8', 'ignore')

        else:
            extracted_text = f"File type '{file_extension}' is not supported."

    except Exception as e:
        return f"An error occurred while processing '{file_name}': {e}"

    return extracted_text


def download_text_as_file(text_content, original_filename):
    """
    Triggers a browser download for the given text content.

    Args:
        text_content (str): The text to be included in the downloadable file.
        original_filename (str): The name of the original uploaded file.
    """
    # Create a new filename for the output file
    base_name, _ = os.path.splitext(original_filename)
    download_filename = f"converted_{base_name}.txt"

    # Use `files.download` to trigger the download in the browser
    with open(download_filename, "w", encoding="utf-8") as f:
        f.write(text_content)
    
    print(f"\nDownloading full text as '{download_filename}'...")
    files.download(download_filename)


# ==============================================================================
# 4. MAIN EXECUTION BLOCK
# ==============================================================================

if __name__ == "__main__":
    # Step 1: Upload the file
    filename, content_stream = upload_file_to_colab()

    # Proceed only if a file was successfully uploaded
    if filename and content_stream:
        # Step 2: Convert the file content to text
        converted_text = convert_file_to_text(filename, content_stream)

        # Step 3: Display the first 1000 characters as a preview
        print("\n" + "="*50)
        print("                CONVERSION PREVIEW (First 1000 characters)")
        print("="*50)
        print(converted_text[:1000])
        print("="*50)

        # Step 4: Give the user the option to download the full text
        if converted_text and not converted_text.startswith("File type"):
            download_text_as_file(converted_text, filename)
        else:
            print("\nDownload skipped due to unsupported file type or conversion error.")
